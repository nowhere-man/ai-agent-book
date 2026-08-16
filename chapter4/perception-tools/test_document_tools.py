"""
Tests for document processing tools.
Uses real files for testing.
"""
import asyncio
import json
import pytest
import tempfile
from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).parent / "src"))

from document_processing_tools import (
    extract_docx_content,
    extract_pptx_content,
    extract_csv_content
)


class TestDOCXExtraction:
    """Tests for DOCX extraction."""
    
    @pytest.mark.asyncio
    async def test_extract_docx(self):
        """Test DOCX extraction."""
        from docx import Document
        
        docx_path = Path(tempfile.mktemp(suffix=".docx"))
        
        # Create DOCX
        doc = Document()
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")
        doc.save(docx_path)
        
        try:
            result = await extract_docx_content(str(docx_path))
            data = json.loads(result.text)
            
            assert data["success"] is True
            message = data["message"]
            assert message["file_type"] == "docx"
            assert message["paragraphs"] >= 2
            assert "First paragraph" in message["text"]
            
            print(f"✅ DOCX extraction: {message['paragraphs']} paragraphs")
        finally:
            docx_path.unlink(missing_ok=True)


class TestPPTXExtraction:
    """Tests for PPTX extraction."""
    
    @pytest.mark.asyncio
    async def test_extract_pptx(self):
        """Test PPTX extraction."""
        from pptx import Presentation
        
        pptx_path = Path(tempfile.mktemp(suffix=".pptx"))
        
        # Create PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide"
        prs.save(pptx_path)
        
        try:
            result = await extract_pptx_content(str(pptx_path))
            data = json.loads(result.text)
            
            assert data["success"] is True
            message = data["message"]
            assert message["file_type"] == "pptx"
            assert message["total_slides"] >= 1
            
            print(f"✅ PPTX extraction: {message['total_slides']} slides")
        finally:
            pptx_path.unlink(missing_ok=True)


class TestCSVExtraction:
    """Tests for CSV extraction."""
    
    @pytest.mark.asyncio
    async def test_extract_csv(self):
        """Test CSV extraction."""
        csv_path = Path(tempfile.mktemp(suffix=".csv"))
        
        # Create CSV
        csv_content = """name,age,city
Alice,30,New York
Bob,25,Los Angeles
Charlie,35,Chicago"""
        csv_path.write_text(csv_content)
        
        try:
            result = await extract_csv_content(str(csv_path))
            data = json.loads(result.text)
            
            assert data["success"] is True
            message = data["message"]
            assert message["file_type"] == "csv"
            assert message["rows"] == 3
            assert message["columns"] == 3
            assert "name" in message["column_names"]
            
            print(f"✅ CSV extraction: {message['rows']} rows, {message['columns']} columns")
        finally:
            csv_path.unlink(missing_ok=True)


if __name__ == "__main__":
    print("=" * 70)
    print("Running Document Processing Tools Tests")
    print("=" * 70)
    print()
    
    pytest.main([__file__, "-v", "-s"])
